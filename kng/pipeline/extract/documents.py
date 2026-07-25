"""Text-format extractors: docx/doc, pdf, pptx, xlsx.

WP1b — Sarvam-first: *every* document is routed through the Sarvam API.
  - PDFs (and images, in media.py): Sarvam Document Intelligence OCR of the
    whole file.
  - Office docs (docx/pptx/xlsx): parsed locally, then the raw text is cleaned
    and structured to Markdown by the Sarvam LLM (`clean_document`).

Sarvam is *primary*; the local parse is the *fallback* — used when Sarvam fails,
there is no key, or the caller passes `sarvam=False` (`--local-only`, dev). Every
paid call is tallied into the shared `calls` counter as ocr / cleanup.

Uniform extractor signature (dispatched by extract/__init__.py):
    fn(path, meta, rel, *, use_ocr, sarvam, use_cleanup, calls) -> list[Segment]
Not every extractor uses every flag; the signature is uniform so dispatch is one
call site.
"""
from __future__ import annotations

import re
import subprocess
from pathlib import Path

from ...models import Locator, Segment, SourceType
from ..normalize import detect_language

# A page/segment with fewer than this many non-space chars is treated as empty.
_MIN_TEXT = 12


def _bump(calls: dict | None, kind: str, n: int = 1) -> None:
    if calls is not None:
        calls[kind] = calls.get(kind, 0) + n


def _seg(meta: dict, rel: str, sid: str, text: str, locator: Locator,
         stype: SourceType | None = None) -> Segment:
    return Segment(
        segment_id=sid,
        source_file=rel,
        source_type=stype or meta["source_type"],
        press_meet_id=meta["press_meet_id"],
        press_meet_title=meta["press_meet_title"],
        date=meta["date"],
        topic=meta["topic"],
        publication=meta["publication"],
        locator=locator,
        text_original=text.strip(),
    )


def _maybe_clean(text: str, *, sarvam: bool, use_cleanup: bool, calls: dict | None) -> str:
    """Route locally-parsed office-doc text through the Sarvam LLM for cleanup.

    Falls back to the raw local text when Sarvam is disabled/unavailable, or when
    the cleaned result looks lossy (much shorter than the input → possible
    truncation/summarisation), which would violate the preserve-all-content rule.
    """
    text = text.strip()
    if not text or not (sarvam and use_cleanup):
        return text
    try:
        from ...providers import get_llm
        cleaned = get_llm().clean_document(text, lang_hint=detect_language(text))
        _bump(calls, "cleanup")
        if cleaned and len(cleaned) >= 0.5 * len(text):
            return cleaned.strip()
        _bump(calls, "cleanup_lossy")       # kept raw text: cleaned looked truncated
    except Exception:
        # Falling back to the raw local parse is by design, but a *silent*
        # fallback hid that every office-doc cleanup call was failing. Tally it
        # so `kng.stats` shows the Sarvam-first requirement is unmet.
        _bump(calls, "cleanup_failed")
    return text


# ── DOCX / DOC ─────────────────────────────────────────────────────────────────
def extract_doc(path: Path, meta: dict, rel: str, *, use_ocr: bool = True,
                sarvam: bool = False, use_cleanup: bool = True,
                calls: dict | None = None) -> list[Segment]:
    if path.suffix.lower() == ".doc":
        text = _legacy_doc_text(path)
    else:
        import docx
        d = docx.Document(str(path))
        parts: list[str] = [p.text for p in d.paragraphs if p.text and p.text.strip()]
        for table in d.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells]
                if any(cells):
                    parts.append(" | ".join(cells))
        text = "\n".join(parts)
    if not text.strip():
        return []
    if _is_legacy_font_mojibake(text):
        raise RuntimeError(
            "text is a legacy non-Unicode Indic font (e.g. Shree-Lipi/Anu), not "
            "convertible without a glyph mapping — skipped rather than indexed")
    text = _maybe_clean(text, sarvam=sarvam, use_cleanup=use_cleanup, calls=calls)
    return [_seg(meta, rel, f"{rel}#doc", text, Locator())]


def _legacy_doc_text(path: Path) -> str:
    """Extract a legacy `.doc`, which in this archive is one of three things.

    The extension lies often enough that the format is decided by magic bytes:
    `*.RTF.doc` files are really RTF, and the rest are OLE2 Word binaries.
    LibreOffice is used when present (most faithful), but both formats have a
    pure-Python path so extraction does not depend on it.
    """
    import shutil
    head = path.open("rb").read(8)
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if soffice:
        return _soffice_text(path, soffice)
    if head.startswith(b"{\\rtf"):
        from striprtf.striprtf import rtf_to_text
        return rtf_to_text(path.read_text(encoding="latin-1"), errors="ignore")
    if head.startswith(b"\xd0\xcf\x11\xe0"):        # OLE2 compound document
        return _ole_word_text(path)
    raise RuntimeError(f"unrecognised legacy .doc container: {head[:4]!r}")


def _soffice_text(path: Path, soffice: str) -> str:
    import tempfile
    with tempfile.TemporaryDirectory() as td:
        subprocess.run([soffice, "--headless", "--convert-to", "txt:Text",
                        "--outdir", td, str(path)], check=True,
                       stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        out = Path(td) / (path.stem + ".txt")
        return out.read_text(encoding="utf-8", errors="replace") if out.exists() else ""


def _ole_word_text(path: Path) -> str:
    """Read the main text span of an OLE2 Word binary without LibreOffice.

    The FIB at offset 0x18 holds fcMin/fcMac, which bound the document text for
    non-complex files. Word's in-band control marks (cell/row/page breaks) become
    newlines instead of being dropped, so words either side stay separated.
    """
    import struct

    import olefile
    with olefile.OleFileIO(str(path)) as ole:
        data = ole.openstream("WordDocument").read()
    fc_min, fc_mac = struct.unpack("<II", data[0x18:0x20])
    if not 0 < fc_min < fc_mac <= len(data):
        raise RuntimeError("OLE2 .doc: FIB text range out of bounds (complex/encrypted?)")
    text = data[fc_min:fc_mac].decode("cp1252", "replace")
    text = re.sub(r"[\r\x07\x0b\x0c]", "\n", text)      # paragraph/cell/page marks
    text = re.sub(r"[\x00-\x08\x0e-\x1f]", "", text)    # residual control bytes
    return re.sub(r"\n{3,}", "\n\n", text)


# Legacy Indic DTP fonts (Shree-Lipi, Anu, APS…) map Telugu glyphs onto cp1252
# byte values. The bytes extract fine but are meaningless as text, and would be
# embedded as if they were Telugu. Detect and reject rather than index garbage.
def _is_legacy_font_mojibake(text: str) -> bool:
    if len(text) < 50:
        return False
    if any(0x0C00 <= ord(c) <= 0x0C7F for c in text):   # real Telugu present
        return False
    high = sum(1 for c in text if 0x80 <= ord(c) <= 0xFF)
    return high / len(text) > 0.15


# ── PDF ────────────────────────────────────────────────────────────────────────
def extract_pdf(path: Path, meta: dict, rel: str, *, use_ocr: bool = True,
                sarvam: bool = False, use_cleanup: bool = True,
                calls: dict | None = None) -> list[Segment]:
    # Sarvam-first: OCR the whole file so scanned + digital PDFs go via the API.
    ocr_err: str | None = None
    if sarvam and use_ocr:
        segs, ocr_err = _pdf_via_sarvam(path, meta, rel, calls)
        if segs:
            return segs
        # fall through to local PyMuPDF text on failure / empty OCR
    segs = _pdf_via_pymupdf(path, meta, rel)
    if not segs and ocr_err:
        # A scanned PDF whose OCR call failed has no text layer either, so both
        # paths come back empty. Reporting success here would save a contentless
        # doc, mark it "done" in the manifest and never retry it.
        raise RuntimeError(f"sarvam OCR failed and PDF has no text layer: {ocr_err}")
    return segs


def _pdf_via_sarvam(path: Path, meta: dict, rel: str,
                    calls: dict | None) -> tuple[list[Segment], str | None]:
    from ...providers import get_ocr
    ocr = get_ocr()
    try:
        pages = ocr.ocr_file(path, languages=["te", "hi", "en"])
    except Exception as e:
        # A long PDF is OCR'd as several billed jobs; if one slice fails the whole
        # file is retried next run rather than saved with pages missing.
        _bump(calls, "ocr", getattr(ocr, "last_call_jobs", 0))
        return [], f"{type(e).__name__}: {e}"
    _bump(calls, "ocr", getattr(ocr, "last_call_jobs", 1) or 1)
    segs: list[Segment] = []
    for page_no, text in pages:
        text = (text or "").strip()
        if len(text) >= _MIN_TEXT:
            segs.append(_seg(meta, rel, f"{rel}#p{page_no}", text,
                             Locator(page=page_no), SourceType.source_doc))
    segs.sort(key=lambda s: s.locator.page or 0)
    return segs, None


def _pdf_via_pymupdf(path: Path, meta: dict, rel: str) -> list[Segment]:
    import fitz
    doc = fitz.open(str(path))
    segs: list[Segment] = []
    for i, page in enumerate(doc, start=1):
        text = page.get_text("text")
        if len(text.strip()) >= _MIN_TEXT:
            segs.append(_seg(meta, rel, f"{rel}#p{i}", text, Locator(page=i)))
    segs.sort(key=lambda s: s.locator.page or 0)
    return segs


# ── PPTX ───────────────────────────────────────────────────────────────────────
def extract_pptx(path: Path, meta: dict, rel: str, *, use_ocr: bool = True,
                 sarvam: bool = False, use_cleanup: bool = True,
                 calls: dict | None = None) -> list[Segment]:
    from pptx import Presentation
    prs = Presentation(str(path))
    segs: list[Segment] = []
    for i, slide in enumerate(prs.slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                lines.append(shape.text_frame.text.strip())
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            note = slide.notes_slide.notes_text_frame.text.strip()
            if note:
                lines.append(f"[notes] {note}")
        text = "\n".join(lines)
        if text.strip():
            text = _maybe_clean(text, sarvam=sarvam, use_cleanup=use_cleanup, calls=calls)
            segs.append(_seg(meta, rel, f"{rel}#s{i}", text, Locator(slide=i)))
    return segs


# ── XLSX ───────────────────────────────────────────────────────────────────────
def extract_xlsx(path: Path, meta: dict, rel: str, *, use_ocr: bool = True,
                 sarvam: bool = False, use_cleanup: bool = True,
                 calls: dict | None = None) -> list[Segment]:
    import pandas as pd
    sheets = pd.read_excel(str(path), sheet_name=None, header=None, dtype=str)
    segs: list[Segment] = []
    for idx, (name, df) in enumerate(sheets.items(), start=1):
        df = df.fillna("")
        rows = ["\t".join(str(c) for c in row) for row in df.values.tolist()[:200]]
        text = f"# Sheet: {name}\n" + "\n".join(r for r in rows if r.strip())
        if text.strip():
            text = _maybe_clean(text, sarvam=sarvam, use_cleanup=use_cleanup, calls=calls)
            segs.append(_seg(meta, rel, f"{rel}#sheet{idx}", text, Locator()))
    return segs
